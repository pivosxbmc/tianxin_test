# -*- coding: utf-8 -*-
from multiprocessing import Process
from pptx import Presentation
from docx import Document
import sys
import importlib
importlib.reload(sys)
wordfile=Document()
# 给定ppt文件所在的路径
filepath='1030.pptx'
pptx = Presentation(filepath)
# 遍历ppt文件的所有幻灯片页
for slide in pptx.slides:
	for shape in slide.shapes:
		if shape.has_text_frame:
			text_frame= shape.text_frame
			for paragraph in text_frame.paragraphs:
				print(paragraph.text)
				#wordfile.add_paragraph(u'%s'%'test')
save_path='text.docx'
wordfile.save(save_path)